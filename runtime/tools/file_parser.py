from __future__ import annotations

import base64
import logging
import tempfile
from dataclasses import dataclass
from pathlib import Path

from config.settings import settings

logger = logging.getLogger("runtime.tools.file_parser")


@dataclass
class FileParseResult:
    text: str
    mime_type: str
    filename: str
    pages: int | None = None
    sheets: list[str] | None = None
    slides: int | None = None
    ocr_used: bool = False


def parse_file(path: Path, mime_type: str, filename: str = "untitled") -> FileParseResult:
    mime_lower = mime_type.lower().strip().split(";")[0].strip()

    if mime_lower == "text/plain":
        text = path.read_text("utf-8", errors="replace")
        return FileParseResult(text=text, mime_type=mime_type, filename=filename)

    if mime_lower == "text/markdown":
        text = path.read_text("utf-8", errors="replace")
        return FileParseResult(text=text, mime_type=mime_type, filename=filename)

    if mime_lower == "application/pdf":
        return _parse_pdf(path, filename)

    if mime_lower == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return _parse_docx(path, filename)

    if mime_lower == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        return _parse_xlsx(path, filename)

    if mime_lower == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return _parse_pptx(path, filename)

    return FileParseResult(
        text=f"[Unsupported file type: {mime_type}]",
        mime_type=mime_type,
        filename=filename,
    )


def parse_base64(data: str, mime_type: str, filename: str = "untitled") -> FileParseResult:
    try:
        raw = base64.b64decode(data)
    except Exception as exc:
        logger.warning("Failed to decode base64 file %s: %s", filename, exc)
        return FileParseResult(
            text=f"[Failed to decode file: {filename}]",
            mime_type=mime_type,
            filename=filename,
        )
    with tempfile.NamedTemporaryFile(delete=False) as f:
        f.write(raw)
        tmp = Path(f.name)
    try:
        return parse_file(tmp, mime_type, filename=filename)
    finally:
        try:
            tmp.unlink(missing_ok=True)
        except OSError:
            pass


def _parse_pdf(path: Path, filename: str) -> FileParseResult:
    import fitz

    with fitz.open(str(path)) as doc:
        texts: list[str] = []
        total_chars = 0
        for page in doc:
            text = page.get_text()
            texts.append(text)
            total_chars += len(text)

        result = FileParseResult(
            text="\n\n".join(texts),
            mime_type="application/pdf",
            filename=filename,
            pages=len(doc),
        )

    if total_chars >= 200:
        return result

    try:
        import pytesseract
    except ImportError:
        result.text = "[OCR not available - pytesseract not installed]"
        return result

    try:
        ocr_texts: list[str] = []
        with fitz.open(str(path)) as doc:
            for page in doc:
                pix = page.get_pixmap(dpi=300)
                img_bytes = pix.tobytes("png")
                ocr_text = pytesseract.image_to_string(img_bytes, lang=settings.tesseract_langs)
                ocr_texts.append(ocr_text.strip())
        combined = "\n\n".join(t for t in ocr_texts if t)
        result = FileParseResult(
            text=combined or "[OCR produced no text]",
            mime_type="application/pdf",
            filename=filename,
            pages=result.pages,
            ocr_used=True,
        )
    except Exception as exc:
        logger.warning("OCR failed for %s: %s", filename, exc)
        result.text = "[OCR failed]"

    return result


def _parse_docx(path: Path, filename: str) -> FileParseResult:
    from docx import Document

    doc = Document(str(path))
    paras = [p.text for p in doc.paragraphs]
    return FileParseResult(
        text="\n\n".join(paras),
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        filename=filename,
    )


def _parse_xlsx(path: Path, filename: str) -> FileParseResult:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    sheet_names: list[str] = []
    for ws in wb.worksheets:
        sheet_names.append(ws.title)
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            rows.append("\t".join(cells))
        parts.append(f"[Sheet: {ws.title}]\n" + "\n".join(rows))
    wb.close()
    return FileParseResult(
        text="\n\n".join(parts),
        mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        filename=filename,
        sheets=sheet_names,
    )


def _parse_pptx(path: Path, filename: str) -> FileParseResult:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    slide_count = 0
    for slide in prs.slides:
        slide_count += 1
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_texts.append("\t".join(cells))
        if slide_texts:
            parts.append(f"[Slide: {slide_count}]\n" + "\n".join(slide_texts))
    return FileParseResult(
        text="\n\n".join(parts),
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=filename,
        slides=slide_count,
    )
