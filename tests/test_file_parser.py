from __future__ import annotations

import base64
import io
import tempfile
from pathlib import Path

import pytest

from runtime.tools.file_parser import FileParseResult, parse_file, parse_base64


# --- helpers ---

def _write_temp(content: bytes | str, suffix: str) -> Path:
    f = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)
    if isinstance(content, str):
        f.write(content.encode("utf-8"))
    else:
        f.write(content)
    f.close()
    return Path(f.name)


def test_parse_txt():
    path = _write_temp("Hello world", ".txt")
    try:
        result = parse_file(path, "text/plain")
        assert isinstance(result, FileParseResult)
        assert "Hello world" in result.text
        assert result.pages is None
    finally:
        path.unlink(missing_ok=True)


def test_parse_md():
    path = _write_temp("# Title\n\nBody text", ".md")
    try:
        result = parse_file(path, "text/markdown")
        assert "# Title" in result.text
        assert "Body text" in result.text
    finally:
        path.unlink(missing_ok=True)


def test_parse_pdf_text():
    import fitz
    doc = fitz.open()
    page = doc.new_page()
    rect = fitz.Rect(50, 50, 550, 800)
    page.insert_textbox(rect, "Hello PDF. " * 200, fontsize=8)
    buf = doc.write()
    doc.close()
    path = _write_temp(buf, ".pdf")
    try:
        result = parse_file(path, "application/pdf")
        assert "Hello PDF" in result.text
        assert isinstance(result.pages, int)
        assert result.pages >= 1
        assert result.ocr_used is False
    finally:
        path.unlink(missing_ok=True)


def test_parse_docx():
    from docx import Document
    doc = Document()
    doc.add_paragraph("Hello DOCX")
    buf = io.BytesIO()
    doc.save(buf)
    path = _write_temp(buf.getvalue(), ".docx")
    try:
        result = parse_file(path, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        assert "Hello DOCX" in result.text
    finally:
        path.unlink(missing_ok=True)


def test_parse_xlsx():
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Name"
    ws["B1"] = "Age"
    ws["A2"] = "Alice"
    ws["B2"] = "30"
    buf = io.BytesIO()
    wb.save(buf)
    path = _write_temp(buf.getvalue(), ".xlsx")
    try:
        result = parse_file(path, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        assert "[Sheet: Sheet1]" in result.text
        assert "Alice" in result.text
        assert "30" in result.text
        assert result.sheets == ["Sheet1"]
    finally:
        path.unlink(missing_ok=True)


def test_parse_pptx():
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = "Hello PPTX"
    buf = io.BytesIO()
    prs.save(buf)
    path = _write_temp(buf.getvalue(), ".pptx")
    try:
        result = parse_file(path, "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        assert "[Slide: 1]" in result.text
        assert "Hello PPTX" in result.text
        assert result.slides == 1
    finally:
        path.unlink(missing_ok=True)


def test_parse_unknown_mime():
    path = _write_temp("whatever", ".bin")
    try:
        result = parse_file(path, "application/octet-stream")
        assert "[Unsupported file type:" in result.text
    finally:
        path.unlink(missing_ok=True)


def test_parse_base64_text():
    b64 = base64.b64encode(b"Hello from base64").decode()
    result = parse_base64(b64, "text/plain", "test.txt")
    assert "Hello from base64" in result.text


def test_parse_base64_pdf():
    import fitz
    doc = fitz.open()
    page = doc.new_page()
    rect = fitz.Rect(50, 50, 550, 800)
    page.insert_textbox(rect, "PDF base64. " * 200, fontsize=8)
    buf = doc.write()
    doc.close()
    b64 = base64.b64encode(buf).decode()
    result = parse_base64(b64, "application/pdf", "doc.pdf")
    assert "PDF base64" in result.text
